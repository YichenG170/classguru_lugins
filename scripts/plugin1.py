"""
AI 文档分析服务（MaterialSvc）
功能：处理用户上传的文档（PPT、PDF、图片、笔记），生成课程画像（CourseProfile JSON）
"""

import json
import os
import base64
import uuid
from typing import Dict, List, Any, Optional
from dataclasses import dataclass, asdict
from pathlib import Path
from datetime import datetime
import pymupdf  # PyMuPDF for PDF processing
from PIL import Image
import openai
from pptx import Presentation


@dataclass
class CourseProfile:
    """课程画像数据结构"""
    main_title: str
    subtitle: str
    introduction: str
    keywords: List[str]  # 30-50个核心术语
    outline: List[str]   # 8-12个主题条目
    abbreviations: Dict[str, str]  # 10-20个缩写映射
    formulas_symbols: List[Dict[str, str]]  # 8-15个变量，包含symbol和description
    proper_nouns_cases: List[str]  # 10-20个专名/案例
    learning_objectives: List[str]  # 3-5条教学目标


@dataclass
class APIRequest:
    """标准化API请求结构"""
    version: str
    request_id: str
    source: Dict[str, Any]
    intent: Dict[str, Any]
    expect: Dict[str, Any]


@dataclass
class APIResponse:
    """标准化API回复结构"""
    request_id: str
    result: Dict[str, Any]
    info: Dict[str, Any]


class MaterialSvc:
    """AI 文档分析服务"""
    
    def __init__(self, openai_api_key: Optional[str] = None):
        """
        初始化服务
        
        Args:
            openai_api_key: OpenAI API密钥。如果为None，则从环境变量OPENAI_API_KEY中获取
            
        Raises:
            ValueError: 无法获取有效的API密钥时抛出
        """
        if openai_api_key is None:
            openai_api_key = os.getenv('OPENAI_API_KEY')
            
        if not openai_api_key:
            raise ValueError(
                "未找到OpenAI API密钥。请通过以下方式之一提供：\n"
                "1. 在初始化时传入 openai_api_key 参数\n"
                "2. 设置环境变量 OPENAI_API_KEY\n"
                "   Windows: $env:OPENAI_API_KEY='your_api_key'\n"
                "   Linux/Mac: export OPENAI_API_KEY='your_api_key'"
            )
            
        self.client = openai.OpenAI(api_key=openai_api_key)
        self.supported_formats = {'.pdf', '.ppt', '.pptx', '.jpg', '.jpeg', '.png', '.txt', '.md'}
        self.version = "1.0.0"
    
    def handle_request(self, request_data: Dict[str, Any]) -> Dict[str, Any]:
        """
        标准化API请求处理入口
        
        Args:
            request_data: 符合API规范的请求字典
            
        Returns:
            Dict[str, Any]: 符合API规范的响应字典
        """
        try:
            # 验证请求格式
            api_request = self._validate_request(request_data)
            
            # 处理业务逻辑
            result = self._process_intent(api_request)
            
            # 构建标准响应
            return self._build_response(api_request.request_id, result, 200, "处理成功")
            
        except ValueError as e:
            return self._build_error_response(
                request_data.get("request_id", str(uuid.uuid4())), 
                400, 
                f"请求参数错误: {str(e)}"
            )
        except Exception as e:
            return self._build_error_response(
                request_data.get("request_id", str(uuid.uuid4())), 
                500, 
                f"服务内部错误: {str(e)}"
            )
    
    def _validate_request(self, request_data: Dict[str, Any]) -> APIRequest:
        """验证API请求格式"""
        required_fields = ["version", "request_id", "source", "intent", "expect"]
        
        for field in required_fields:
            if field not in request_data:
                raise ValueError(f"缺少必需字段: {field}")
        
        # 验证版本兼容性
        if request_data["version"] != self.version:
            raise ValueError(f"版本不兼容: 期望 {self.version}, 收到 {request_data['version']}")
        
        return APIRequest(**request_data)
    
    def _process_intent(self, api_request: APIRequest) -> Dict[str, Any]:
        """处理具体的业务意图"""
        intent = api_request.intent
        action = intent.get("action")
        
        if action == "analyze":
            # 统一分析功能：支持文件路径或直接内容
            file_path = intent.get("file_path")
            content = intent.get("content")
            
            if file_path:
                # 处理文件
                course_profile, openai_file_id = self._analyze_document(file_path)
            elif content:
                # 处理直接内容
                content_type = intent.get("content_type", "text")
                course_profile_data = self._analyze_with_openai(content, content_type)
                course_profile = CourseProfile(**course_profile_data)
                openai_file_id = None
            else:
                raise ValueError("analyze操作需要file_path或content参数")
            
            # 根据expect字段决定返回格式
            expect = api_request.expect
            result = self._format_result(course_profile, expect)
            
            # 添加OpenAI file ID
            if openai_file_id:
                result["openai_file_id"] = openai_file_id
            
            return result
            
        else:
            raise ValueError(f"不支持的操作: {action}")
    
    def _format_result(self, course_profile: CourseProfile, expect: Dict[str, Any]) -> Dict[str, Any]:
        """根据expect字段格式化返回结果"""
        result = {}
        
        # 默认返回完整的课程画像
        if expect.get("format") == "full" or not expect.get("fields"):
            result = asdict(course_profile)
        else:
            # 只返回指定字段
            fields = expect.get("fields", [])
            profile_dict = asdict(course_profile)
            for field in fields:
                if field in profile_dict:
                    result[field] = profile_dict[field]
        
        # 添加额外的元数据
        if expect.get("include_metadata", False):
            result["metadata"] = {
                "generated_at": datetime.now().isoformat(),
                "service_version": self.version,
                "total_keywords": len(course_profile.keywords),
                "total_outline_items": len(course_profile.outline)
            }
        
        return result
    
    def _build_response(self, request_id: str, result: Dict[str, Any], 
                       status_code: int, description: str) -> Dict[str, Any]:
        """构建标准API响应"""
        return {
            "request_id": request_id,
            "result": result,
            "info": {
                "status_code": status_code,
                "description": description
            }
        }
    
    def _build_error_response(self, request_id: str, status_code: int, 
                             description: str) -> Dict[str, Any]:
        """构建错误响应"""
        return {
            "request_id": request_id,
            "result": {},
            "info": {
                "status_code": status_code,
                "description": description
            }
        }
    
    def _analyze_document(self, file_path: str) -> tuple[CourseProfile, Optional[str]]:
        """
        分析文档并生成课程画像，同时上传到OpenAI获取file ID
        
        Args:
            file_path: 文档文件路径
            
        Returns:
            tuple: (CourseProfile对象, OpenAI file ID)
            
        Raises:
            ValueError: 不支持的文件格式
            FileNotFoundError: 文件不存在
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        file_ext = Path(file_path).suffix.lower()
        if file_ext not in self.supported_formats:
            raise ValueError(f"不支持的文件格式: {file_ext}")
        
        # 尝试上传文件到OpenAI获取file ID
        openai_file_id = None
        try:
            with open(file_path, "rb") as f:
                file_obj = self.client.files.create(
                    file=f,
                    purpose="assistants"
                )
                openai_file_id = file_obj.id
        except Exception as e:
            print(f"警告: 上传文件到OpenAI失败: {e}")
        
        # 提取文档内容
        content = self._extract_content(file_path, file_ext)
        
        # 调用OpenAI分析内容
        course_profile_data = self._analyze_with_openai(content, file_ext)
        
        # 创建CourseProfile对象
        course_profile = CourseProfile(**course_profile_data)
        
        return course_profile, openai_file_id
    
    def _extract_content(self, file_path: str, file_ext: str) -> str:
        """
        根据文件类型提取内容
        
        Args:
            file_path: 文件路径
            file_ext: 文件扩展名
            
        Returns:
            str: 提取的文本内容
        """
        try:
            if file_ext == '.pdf':
                return self._extract_pdf_content(file_path)
            elif file_ext in ['.ppt', '.pptx']:
                return self._extract_ppt_content(file_path)
            elif file_ext in ['.jpg', '.jpeg', '.png']:
                return self._extract_image_content(file_path)
            elif file_ext in ['.txt', '.md']:
                return self._extract_text_content(file_path)
            else:
                raise ValueError(f"暂不支持的文件格式: {file_ext}")
        except Exception as e:
            raise Exception(f"内容提取失败: {str(e)}")
    
    def _extract_pdf_content(self, file_path: str) -> str:
        """提取PDF内容"""
        content = []
        doc = pymupdf.Document(file_path)
        
        for page_num in range(doc.page_count):
            page = doc[page_num]
            content.append(page.get_text())
        
        doc.close()
        return "\n\n".join(content)
    
    def _extract_ppt_content(self, file_path: str) -> str:
        """提取PPT内容"""
        content = []
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = [f"=== 幻灯片 {slide_num} ==="]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            
            content.append("\n".join(slide_content))
        
        return "\n\n".join(content)
    
    def _extract_image_content(self, file_path: str) -> str:
        """
        提取图片内容（返回base64编码用于OpenAI Vision API）
        """
        with open(file_path, "rb") as image_file:
            encoded_image = base64.b64encode(image_file.read()).decode('utf-8')
        return f"[IMAGE_BASE64]{encoded_image}"
    
    def _extract_text_content(self, file_path: str) -> str:
        """提取文本文件内容"""
        encodings = ['utf-8', 'gbk', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        
        raise Exception("无法解码文本文件")
    
    def _analyze_with_openai(self, content: str, file_ext: str) -> Dict[str, Any]:
        """
        使用OpenAI分析内容并生成课程画像
        
        Args:
            content: 文档内容
            file_ext: 文件扩展名
            
        Returns:
            Dict: 课程画像数据
        """
        # 构建系统提示词
        system_prompt = self._build_system_prompt()
        
        # 构建用户消息
        if content.startswith("[IMAGE_BASE64]"):
            # 处理图片
            base64_image = content[14:]  # 去掉[IMAGE_BASE64]前缀
            messages = [
                {"role": "system", "content": system_prompt},
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "请分析这张图片中的教学内容，并生成课程画像JSON。"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ]
            model = "gpt-4o-mini"
        else:
            # 处理文本内容
            user_content = f"请分析以下{file_ext}文档内容，生成课程画像JSON：\n\n{content}"
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_content}
            ]
            model = "gpt-4o-mini"
        
        try:
            response = self.client.chat.completions.create(
                model=model,
                messages=messages,
                temperature=0.3,
                max_tokens=2000
            )
            
            # 解析响应
            response_content = response.choices[0].message.content
            
            # 提取JSON内容
            json_start = response_content.find('{')
            json_end = response_content.rfind('}') + 1
            
            if json_start == -1 or json_end == 0:
                raise Exception("OpenAI响应中未找到有效的JSON")
            
            json_content = response_content[json_start:json_end]
            return json.loads(json_content)
            
        except Exception as e:
            raise Exception(f"OpenAI分析失败: {str(e)}")
    
    def _build_system_prompt(self) -> str:
        """构建系统提示词"""
        return """你是一个专业的教育内容分析师。请根据提供的文档内容，生成标准的课程画像JSON。

要求：
1. 严格按照以下JSON Schema输出
2. 所有内容必须来自原始文档的显式信息，不得编造
3. 如果某些信息在文档中不存在，使用空数组[]或空字符串""

JSON Schema:
{
    "main_title": "主要标题",
    "subtitle": "副标题", 
    "introduction": "课程简介（简要说明课程内容）",
    "keywords": ["关键词1", "关键词2", ...],  // 30-50个核心术语
    "outline": ["主题1", "主题2", ...],  // 8-12个主题条目
    "abbreviations": {"缩写1": "全称1", "缩写2": "全称2", ...},  // 10-20个缩写映射
    "formulas_symbols": [
        {"symbol": "符号", "description": "描述"},
        ...
    ],  // 8-15个变量
    "proper_nouns_cases": ["专名1", "案例1", ...],  // 10-20个专名/案例
    "learning_objectives": ["目标1", "目标2", ...]  // 3-5条教学目标
}

请直接输出符合此格式的JSON，不要添加额外的说明文字。"""

    def save_course_profile(self, course_profile: CourseProfile, output_path: str) -> None:
        """
        保存课程画像到JSON文件
        
        Args:
            course_profile: 课程画像对象
            output_path: 输出文件路径
        """
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(asdict(course_profile), f, ensure_ascii=False, indent=2)
            print(f"课程画像已保存到: {output_path}")
        except Exception as e:
            raise Exception(f"保存文件失败: {str(e)}")
    
    def print_course_profile(self, course_profile: CourseProfile) -> None:
        """打印课程画像信息"""
        print("=== 课程画像 ===")
        print(f"主标题: {course_profile.main_title}")
        print(f"副标题: {course_profile.subtitle}")
        print(f"简介: {course_profile.introduction}")
        print(f"关键词数量: {len(course_profile.keywords)}")
        print(f"大纲条目数: {len(course_profile.outline)}")
        print(f"缩写数量: {len(course_profile.abbreviations)}")
        print(f"公式符号数: {len(course_profile.formulas_symbols)}")
        print(f"专名案例数: {len(course_profile.proper_nouns_cases)}")
        print(f"学习目标数: {len(course_profile.learning_objectives)}")


def main():
    """主函数示例"""
    try:
        # 初始化服务
        material_svc = MaterialSvc()
        
        # 测试请求：分析PDF文档
        request_data_1 = {
            "version": "1.0.0",
            "request_id": str(uuid.uuid4()),
            "source": {
                "timestamp": datetime.now().isoformat(),
                "page": "/class/123/material",
                "app": {"name": "classguru-web", "version": "1.4.2"},
                "locale": "zh-CN",
                "timezone": "Asia/Shanghai"
            },
            "intent": {
                "action": "analyze",
                "file_path": "cgv1-2 2024.pdf"
            },
            "expect": {
                "format": "full",
                "include_metadata": True
            }
        }
        
        print("=== 请求体1 ===")
        print(json.dumps(request_data_1, indent=2, ensure_ascii=False))
        
        response_1 = material_svc.handle_request(request_data_1)
        
        print("\n=== 回复体1 ===")
        print(json.dumps(response_1, indent=2, ensure_ascii=False))
    except Exception as e:
        print(f"❌ 错误: {str(e)}")


def demo_api_usage():
    """API使用示例"""
    request_example = {
        "version": "1.0.0",
        "request_id": "01ARZ3NDEKTSV4RRFFQ69G5FAV",
        "source": {
            "timestamp": "2025-09-26T12:00:00Z",
            "page": "/class/123/report",
            "app": {"name": "classguru-web", "version": "1.4.2"},
            "locale": "zh-CN",
            "timezone": "Asia/Shanghai"
        },
        "intent": {
            "action": "analyze",
            "file_path": "/path/to/document.pdf"
        },
        "expect": {
            "format": "full",
            "include_metadata": True
        }
    }
    
    print("=== API请求示例模板 ===")
    print(json.dumps(request_example, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()